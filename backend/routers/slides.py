from fastapi import APIRouter, HTTPException
from pydantic import BaseModel, Field
from typing import Optional, List
from groq import Groq
from config import get_config
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import uuid
import json

router = APIRouter(prefix="/api/slides", tags=["slides"])

config = get_config()


def get_groq_client():
    """Create the Groq client lazily so backend startup does not crash."""
    if not config.GROQ_API_KEY:
        return None

    try:
        return Groq(api_key=config.GROQ_API_KEY)
    except Exception as error:
        print(f"Groq client setup failed, using slide fallback: {str(error)}")
        return None

# Theme colors
THEMES = {
    "professional": {
        "title": RGBColor(31, 78, 121),  # Dark blue
        "text": RGBColor(51, 51, 51),    # Dark gray
        "accent": RGBColor(192, 0, 0),   # Red
        "bg": RGBColor(255, 255, 255)    # White
    },
    "modern": {
        "title": RGBColor(0, 102, 204),  # Bright blue
        "text": RGBColor(64, 64, 64),    # Medium gray
        "accent": RGBColor(255, 153, 0), # Orange
        "bg": RGBColor(245, 245, 245)    # Light gray
    },
    "creative": {
        "title": RGBColor(153, 51, 153), # Purple
        "text": RGBColor(51, 51, 51),    # Dark gray
        "accent": RGBColor(51, 204, 204),# Teal
        "bg": RGBColor(240, 240, 250)    # Light blue
    }
}


class ScriptSegment(BaseModel):
    slide_number: int
    title: Optional[str] = None
    text: Optional[str] = None
    content: Optional[str] = None
    word_count: Optional[int] = None
    duration_secs: Optional[int] = None


class SlidesRequest(BaseModel):
    topic: str
    script_segments: List[ScriptSegment]
    theme: Optional[str] = "professional"
    num_slides: Optional[int] = None


class SlidesResponse(BaseModel):
    slides_id: str
    topic: str
    num_slides: int
    file_path: str
    file_name: str
    theme: str
    slides: List[dict] = Field(default_factory=list)
    total_slides: int


def get_segment_text(segment: ScriptSegment) -> str:
    """Accept either `text` or `content` so older and newer frontend code both work."""
    return (segment.text or segment.content or "").strip()


def generate_bullet_points(segment_text: str) -> List[str]:
    """Use Groq to generate bullet points from segment text"""
    try:
        client = get_groq_client()
        if not client:
            raise Exception("Groq client unavailable")

        prompt = f"""Convert this text into 3-4 concise bullet points. Each bullet should be max 8 words.

Text: {segment_text}

Return ONLY the bullet points as a JSON array of strings, nothing else. Example format:
["Bullet point one here", "Second bullet point", "Third point"]"""

        message = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            max_tokens=256,
            messages=[
                {
                    "role": "user",
                    "content": prompt
                }
            ]
        )
        
        response_text = message.choices[0].message.content.strip()
        
        # Parse JSON array
        try:
            bullets = json.loads(response_text)
            return bullets[:4]  # Max 4 bullets
        except json.JSONDecodeError:
            # Fallback: split by newlines or dashes
            bullets = [line.strip().lstrip('- •').strip() 
                      for line in response_text.split('\n') 
                      if line.strip()]
            return bullets[:4]
            
    except Exception as e:
        print(f"Error generating bullet points: {str(e)}")
        # Fallback: create simple bullet points from text
        words = segment_text.split()
        if len(words) > 20:
            return [" ".join(words[:10]), " ".join(words[10:20])]
        return [segment_text]


def create_powerpoint(topic: str, segments: List[ScriptSegment], theme: str = "professional") -> str:
    """Create a PowerPoint presentation from script segments"""
    try:
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        colors = THEMES.get(theme, THEMES["professional"])
        
        # The first slide gives the exported deck a simple cover page.
        title_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Add title background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors["title"]
        
        # Add title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = topic
        title_frame.paragraphs[0].font.size = Pt(60)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"{len(segments)} Slides"
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Each script segment becomes one content slide.
        for segment in segments:
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add white background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = colors["bg"]
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_text = segment.title or f"Slide {segment.slide_number}"
            title_frame.text = title_text
            title_frame.paragraphs[0].font.size = Pt(44)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = colors["title"]
            
            # Add accent line under title
            line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.6), Inches(9), Inches(0))
            line.line.color.rgb = colors["accent"]
            line.line.width = Pt(3)
            
            # Generate and add bullet points
            segment_text = get_segment_text(segment)
            bullets = generate_bullet_points(segment_text)
            
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(4.8))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = bullet
                p.level = 0
                p.font.size = Pt(28)
                p.font.color.rgb = colors["text"]
                p.space_before = Pt(12)
                p.space_after = Pt(12)
        
        # Save presentation
        output_dir = config.OUTPUT_DIR
        os.makedirs(output_dir, exist_ok=True)
        
        slides_id = f"slides_{uuid.uuid4().hex[:8]}"
        file_name = f"{slides_id}.pptx"
        file_path = os.path.join(output_dir, file_name)
        
        prs.save(file_path)
        
        return file_path, file_name, slides_id
        
    except Exception as e:
        raise Exception(f"Failed to create PowerPoint: {str(e)}")


@router.post("/generate", response_model=SlidesResponse)
async def generate_slides(request: SlidesRequest):
    """Generate slides from script segments and create PowerPoint file"""
    try:
        if not request.script_segments:
            raise HTTPException(
                status_code=400,
                detail="script_segments cannot be empty"
            )

        normalized_segments = []
        slide_previews = []

        for index, segment in enumerate(request.script_segments, start=1):
            segment_text = get_segment_text(segment)
            if not segment_text:
                continue

            # Normalize the slide data once so the rest of the function stays simple.
            normalized_segment = ScriptSegment(
                slide_number=segment.slide_number or index,
                title=segment.title,
                text=segment_text,
                word_count=segment.word_count,
                duration_secs=segment.duration_secs
            )
            normalized_segments.append(normalized_segment)
            slide_previews.append({
                "title": normalized_segment.title or f"Slide {normalized_segment.slide_number}",
                "description": segment_text[:180]
            })

        if not normalized_segments:
            raise HTTPException(
                status_code=400,
                detail="script_segments must include text or content"
            )
        
        # Validate theme
        if request.theme not in THEMES:
            request.theme = "professional"
        
        # Create PowerPoint presentation
        file_path, file_name, slides_id = create_powerpoint(
            request.topic,
            normalized_segments,
            request.theme
        )
        
        # Verify file was created
        if not os.path.exists(file_path):
            raise Exception("PowerPoint file was not created successfully")
        
        return SlidesResponse(
            slides_id=slides_id,
            topic=request.topic,
            num_slides=len(normalized_segments) + 1,  # +1 for title slide
            file_path=file_path,
            file_name=file_name,
            theme=request.theme,
            slides=slide_previews,
            total_slides=len(normalized_segments)
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Slide generation failed: {str(e)}")


@router.get("/slides/{slides_id}", response_model=SlidesResponse)
async def get_slides(slides_id: str):
    """Get slides by ID"""
    try:
        # TODO: Fetch slides from database
        return {
            "slides_id": slides_id,
            "script_id": "script_123",
            "slides": [],
            "total_slides": 0
        }
    except Exception as e:
        raise HTTPException(status_code=404, detail="Slides not found")


@router.put("/slides/{slides_id}")
async def update_slides(slides_id: str, request: SlidesRequest):
    """Update slides"""
    try:
        # TODO: Update slides in database
        return {"message": f"Slides {slides_id} updated successfully"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.delete("/slides/{slides_id}")
async def delete_slides(slides_id: str):
    """Delete slides"""
    try:
        # TODO: Delete slides from database
        return {"message": f"Slides {slides_id} deleted successfully"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
