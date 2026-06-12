from fastapi import APIRouter, BackgroundTasks, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import Dict, List, Optional
from config import get_config
from storage import LIBRARY_FILE, SHARES_FILE, load_json_list, save_json_list
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
import os
import shutil
import subprocess
import tempfile
import time
import uuid

router = APIRouter(prefix="/api/video", tags=["video"])
config = get_config()

FFMPEG_PATH = config.FFMPEG_PATH
VIDEO_JOBS: Dict[str, dict] = {}


class AvatarProfile(BaseModel):
    presenter_name: Optional[str] = "Ava"
    appearance: Optional[str] = "mentor"
    style: Optional[str] = "friendly"
    voice_type: Optional[str] = "clear"
    speaking_pace: Optional[str] = "normal"
    provider: Optional[str] = "local"


class AccessControl(BaseModel):
    creator_name: Optional[str] = "Demo Creator"
    creator_role: Optional[str] = "teacher"
    institution: Optional[str] = ""
    access_level: Optional[str] = "private"
    access_code: Optional[str] = ""
    allow_download: Optional[bool] = True


class VideoRequest(BaseModel):
    voice_id: str
    slides_id: str
    add_subtitles: Optional[bool] = False
    project_title: Optional[str] = "Project Video"
    render_mode: Optional[str] = "no-face"
    avatar_profile: Optional[AvatarProfile] = None
    access_control: Optional[AccessControl] = None


class VideoResponse(BaseModel):
    video_id: str
    voice_id: str
    slides_id: str
    file_path: str
    file_name: str
    duration_secs: Optional[int] = None
    status: str
    creator_name: Optional[str] = None
    creator_role: Optional[str] = None
    access_level: Optional[str] = None
    share_token: Optional[str] = None
    share_url: Optional[str] = None


class PreviewResponse(BaseModel):
    preview_id: str
    file_path: str
    file_name: str
    status: str


class RenderJobResponse(BaseModel):
    job_id: str
    status: str
    progress: int
    video_id: Optional[str] = None
    file_path: Optional[str] = None
    file_name: Optional[str] = None
    error: Optional[str] = None
    share_token: Optional[str] = None
    share_url: Optional[str] = None


class ShareRequest(BaseModel):
    video_id: str
    creator_name: Optional[str] = "Demo Creator"
    creator_role: Optional[str] = "teacher"
    institution: Optional[str] = ""
    access_level: Optional[str] = "private"
    access_code: Optional[str] = ""
    allow_download: Optional[bool] = True


class ShareResponse(BaseModel):
    video_id: str
    share_token: str
    share_url: str
    access_level: str
    access_code: Optional[str] = None


def get_font(size: int):
    try:
        return ImageFont.truetype("arial.ttf", size)
    except OSError:
        return ImageFont.load_default()


def ensure_ffmpeg_exists() -> None:
    if not os.path.exists(FFMPEG_PATH):
        raise Exception(f"FFmpeg not found at: {FFMPEG_PATH}")


def get_api_origin() -> str:
    origins = [origin.strip() for origin in config.ALLOWED_ORIGINS.split(",") if origin.strip()]
    https_origins = [origin for origin in origins if origin.startswith("https://")]
    if https_origins:
        return https_origins[0]
    if origins:
        return origins[0]
    return "http://localhost:5173"


def update_job(job_id: str, **updates) -> None:
    VIDEO_JOBS[job_id] = {**VIDEO_JOBS.get(job_id, {}), **updates}


def get_audio_files(voice_id: str) -> List[str]:
    audio_files = []
    if os.path.exists(config.TEMP_DIR):
        for file_name in os.listdir(config.TEMP_DIR):
            if file_name.startswith(voice_id) and file_name.endswith(".mp3"):
                audio_files.append(os.path.join(config.TEMP_DIR, file_name))

    audio_files.sort(key=lambda path: int(os.path.basename(path).split("_")[-1].split(".")[0]))
    return audio_files


def render_slide_image(slide, image_path: str, render_mode: str, avatar_profile: Optional[AvatarProfile]) -> None:
    width, height = 1920, 1080
    image = Image.new("RGB", (width, height), color="#fbfcff")
    draw = ImageDraw.Draw(image)

    title_font = get_font(72)
    text_font = get_font(38)
    small_font = get_font(24)
    medium_font = get_font(30)

    draw.rounded_rectangle((70, 70, 1850, 1010), radius=40, fill="#ffffff", outline="#d9e1ec", width=4)
    draw.rounded_rectangle((95, 95, 1825, 220), radius=30, fill="#eef4ff")

    slide_text = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            slide_text.append(shape.text.strip())

    title = slide_text[0][:55] if slide_text else "Slide"
    draw.text((300, 155), title, fill="#14213d", anchor="lm", font=title_font)
    draw.line((130, 220, 1790, 220), fill="#f2c66f", width=6)

    y_position = 290
    for text in slide_text[1:]:
        words = text.split()
        current_line = ""
        for word in words:
            next_line = f"{current_line} {word}".strip()
            if draw.textlength(next_line, font=text_font) < 1050:
                current_line = next_line
            else:
                draw.text((150, y_position), f"- {current_line}", fill="#29364f", font=text_font)
                y_position += 58
                current_line = word
        if current_line:
            draw.text((150, y_position), f"- {current_line}", fill="#29364f", font=text_font)
            y_position += 58
        y_position += 15

    if render_mode == "avatar":
        presenter_name = avatar_profile.presenter_name if avatar_profile else "Ava"
        appearance = avatar_profile.appearance if avatar_profile else "mentor"
        style = avatar_profile.style if avatar_profile else "friendly"
        panel_x1, panel_y1, panel_x2, panel_y2 = 1300, 260, 1785, 950
        draw.rounded_rectangle((panel_x1, panel_y1, panel_x2, panel_y2), radius=34, fill="#f4f7fb", outline="#cad4e3", width=4)
        draw.rounded_rectangle((1330, 285, 1755, 915), radius=28, fill="#dfe9f7")

        # Shadow behind the avatar
        draw.ellipse((1445, 360, 1705, 625), fill="#b7c8e6")
        # Head and neck
        draw.ellipse((1460, 330, 1690, 560), fill="#f0c9a5", outline="#8b6a54", width=3)
        draw.rectangle((1550, 540, 1600, 610), fill="#e3b792")
        # Hair
        draw.pieslice((1440, 300, 1710, 580), start=180, end=360, fill="#48342a")
        draw.ellipse((1458, 332, 1692, 470), fill="#48342a")
        draw.ellipse((1478, 360, 1672, 560), fill="#f0c9a5")
        # Eyes and mouth
        draw.ellipse((1510, 425, 1540, 438), fill="#25344c")
        draw.ellipse((1610, 425, 1640, 438), fill="#25344c")
        draw.arc((1550, 470, 1600, 510), start=10, end=170, fill="#a25b5b", width=3)
        # Body and jacket
        draw.rounded_rectangle((1435, 590, 1715, 880), radius=36, fill="#284a7c")
        draw.polygon([(1575, 595), (1498, 760), (1652, 760)], fill="#ffffff")
        draw.rectangle((1559, 595, 1591, 760), fill="#d0d9e8")
        draw.line((1495, 760, 1575, 650), fill="#1d355d", width=5)
        draw.line((1655, 760, 1575, 650), fill="#1d355d", width=5)

        draw.text((1545, 835), presenter_name, fill="#ffffff", anchor="mm", font=medium_font)
        draw.text((1545, 875), appearance.title(), fill="#dbe8ff", anchor="mm", font=small_font)
        draw.text((1545, 905), style.title(), fill="#dbe8ff", anchor="mm", font=small_font)
    else:
        draw.rounded_rectangle((1380, 760, 1765, 920), radius=28, fill="#f8f4ea", outline="#e3dac3", width=3)
        draw.text((1415, 805), "No-Face Mode", fill="#8a5a06", anchor="lm", font=medium_font)
        draw.text((1415, 855), "Slides + narration only", fill="#5b667a", anchor="lm", font=small_font)
        draw.text((1415, 895), "Ideal for lectures and demos", fill="#5b667a", anchor="lm", font=small_font)

    image.save(image_path)


def extract_slides_as_images(
    pptx_path: str,
    output_dir: str,
    render_mode: str,
    avatar_profile: Optional[AvatarProfile]
) -> List[str]:
    try:
        presentation = Presentation(pptx_path)
        image_paths = []

        for index, slide in enumerate(presentation.slides, start=1):
            image_path = os.path.join(output_dir, f"slide_{index}.png")
            render_slide_image(slide, image_path, render_mode, avatar_profile)
            image_paths.append(image_path)

        return image_paths
    except Exception as error:
        raise Exception(f"Failed to extract slides: {error}")


def create_video_clip(image_path: str, audio_path: str, output_path: str, duration: int) -> None:
    ensure_ffmpeg_exists()
    command = [
        FFMPEG_PATH,
        "-loop", "1",
        "-i", image_path,
        "-i", audio_path,
        "-c:v", "libx264",
        "-c:a", "aac",
        "-shortest",
        "-pix_fmt", "yuv420p",
        "-t", str(duration),
        "-y",
        output_path
    ]
    result = subprocess.run(command, capture_output=True, text=True, timeout=120)
    if result.returncode != 0:
        raise Exception(f"FFmpeg failed: {result.stderr}")


def concatenate_videos(clip_paths: List[str], output_path: str) -> None:
    ensure_ffmpeg_exists()
    if len(clip_paths) == 1:
        shutil.copyfile(clip_paths[0], output_path)
        return

    command = [FFMPEG_PATH]
    filter_inputs = []

    for index, clip_path in enumerate(clip_paths):
        command.extend(["-i", clip_path])
        filter_inputs.append(f"[{index}:v:0][{index}:a:0]")

    filter_complex = "".join(filter_inputs) + f"concat=n={len(clip_paths)}:v=1:a=1[outv][outa]"

    command.extend([
        "-filter_complex", filter_complex,
        "-map", "[outv]",
        "-map", "[outa]",
        "-c:v", "libx264",
        "-c:a", "aac",
        "-pix_fmt", "yuv420p",
        "-y",
        output_path
    ])

    result = subprocess.run(command, capture_output=True, text=True, timeout=240)
    if result.returncode != 0:
        raise Exception(f"FFmpeg concat failed: {result.stderr}")


def estimate_duration(audio_path: str) -> int:
    audio_size = os.path.getsize(audio_path)
    return max(5, min(30, audio_size // 10000))


def create_share_record(video_id: str, access_control: AccessControl) -> dict:
    shares = load_json_list(SHARES_FILE)
    share_token = f"share_{uuid.uuid4().hex[:10]}"
    share_record = {
        "video_id": video_id,
        "share_token": share_token,
        "access_level": access_control.access_level,
        "access_code": access_control.access_code or "",
        "allow_download": access_control.allow_download,
        "creator_name": access_control.creator_name,
        "creator_role": access_control.creator_role,
        "institution": access_control.institution
    }
    shares.append(share_record)
    save_json_list(SHARES_FILE, shares)
    return share_record


def save_video_record(record: dict) -> None:
    library = load_json_list(LIBRARY_FILE)
    library = [item for item in library if item.get("video_id") != record["video_id"]]
    library.append(record)
    save_json_list(LIBRARY_FILE, library)


def build_video(request: VideoRequest, output_name: str, preview_mode: bool = False) -> dict:
    ensure_ffmpeg_exists()
    temp_dir = os.path.join(config.TEMP_DIR, f"video_{uuid.uuid4().hex[:8]}")
    os.makedirs(temp_dir, exist_ok=True)

    pptx_path = os.path.join(config.OUTPUT_DIR, f"{request.slides_id}.pptx")
    if not os.path.exists(pptx_path):
        raise HTTPException(status_code=404, detail=f"PPTX file not found: {request.slides_id}.pptx")

    audio_files = get_audio_files(request.voice_id)
    if not audio_files:
        raise HTTPException(status_code=404, detail=f"No audio files found for voice_id: {request.voice_id}")

    slide_images = extract_slides_as_images(
        pptx_path,
        temp_dir,
        request.render_mode or "no-face",
        request.avatar_profile
    )

    if len(slide_images) > 1:
        slide_images = slide_images[1:]

    if preview_mode:
        slide_images = slide_images[:1]
        audio_files = audio_files[:1]

    if not slide_images:
        raise HTTPException(status_code=400, detail="No slide images were created for rendering")

    if not audio_files:
        raise HTTPException(status_code=400, detail="No audio files were found for rendering")

    clip_paths = []
    total_duration = 0

    for index, (image_path, audio_path) in enumerate(zip(slide_images, audio_files), start=1):
        clip_path = os.path.join(temp_dir, f"clip_{index}.mp4")
        duration = estimate_duration(audio_path)
        create_video_clip(image_path, audio_path, clip_path, duration)
        clip_paths.append(clip_path)
        total_duration += duration

    if not clip_paths:
        raise HTTPException(
            status_code=400,
            detail="Slides and voice segments could not be matched. Generate slides and voice again."
        )

    final_video_path = os.path.join(config.OUTPUT_DIR, output_name)
    concatenate_videos(clip_paths, final_video_path)
    shutil.rmtree(temp_dir, ignore_errors=True)

    return {
        "file_path": f"outputs/{output_name}",
        "file_name": output_name,
        "duration_secs": total_duration
    }


def process_render_job(job_id: str, request: VideoRequest) -> None:
    try:
        update_job(job_id, status="processing", progress=15)
        time.sleep(1)

        video_id = f"video_{uuid.uuid4().hex[:8]}"
        output = build_video(request, f"{video_id}.mp4")
        update_job(job_id, progress=85)

        access_control = request.access_control or AccessControl()
        share_record = create_share_record(video_id, access_control)
        video_record = {
            "video_id": video_id,
            "voice_id": request.voice_id,
            "slides_id": request.slides_id,
            "file_path": output["file_path"],
            "file_name": output["file_name"],
            "duration_secs": output["duration_secs"],
            "status": "completed",
            "project_title": request.project_title,
            "render_mode": request.render_mode,
            "creator_name": access_control.creator_name,
            "creator_role": access_control.creator_role,
            "institution": access_control.institution,
            "access_level": access_control.access_level,
            "share_token": share_record["share_token"],
            "share_url": f"{get_api_origin()}/share/{share_record['share_token']}"
        }
        save_video_record(video_record)

        update_job(
            job_id,
            status="completed",
            progress=100,
            video_id=video_id,
            file_path=output["file_path"],
            file_name=output["file_name"],
            share_token=share_record["share_token"],
            share_url=video_record["share_url"]
        )
    except HTTPException as error:
        print(f"Render job {job_id} failed with HTTP error: {error.detail}")
        update_job(job_id, status="failed", error=error.detail)
    except Exception as error:
        print(f"Render job {job_id} failed with error: {str(error)}")
        update_job(job_id, status="failed", error=str(error))


@router.post("/preview", response_model=PreviewResponse)
async def generate_preview(request: VideoRequest):
    try:
        preview_id = f"preview_{uuid.uuid4().hex[:8]}"
        output = build_video(request, f"{preview_id}.mp4", preview_mode=True)
        return PreviewResponse(
            preview_id=preview_id,
            file_path=output["file_path"],
            file_name=output["file_name"],
            status="completed"
        )
    except HTTPException:
        raise
    except Exception as error:
        print(f"Preview generation failed: {str(error)}")
        raise HTTPException(status_code=500, detail=f"Preview generation failed: {error}")


@router.post("/jobs", response_model=RenderJobResponse)
async def queue_render(request: VideoRequest, background_tasks: BackgroundTasks):
    job_id = f"job_{uuid.uuid4().hex[:8]}"
    update_job(job_id, job_id=job_id, status="queued", progress=5)
    background_tasks.add_task(process_render_job, job_id, request)
    return RenderJobResponse(job_id=job_id, status="queued", progress=5)


@router.get("/jobs/{job_id}", response_model=RenderJobResponse)
async def get_render_job(job_id: str):
    job = VIDEO_JOBS.get(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Render job not found")
    return RenderJobResponse(**job)


@router.post("/generate", response_model=VideoResponse)
async def generate_video(request: VideoRequest):
    try:
        video_id = f"video_{uuid.uuid4().hex[:8]}"
        output = build_video(request, f"{video_id}.mp4")
        access_control = request.access_control or AccessControl()
        share_record = create_share_record(video_id, access_control)
        record = {
            "video_id": video_id,
            "voice_id": request.voice_id,
            "slides_id": request.slides_id,
            "file_path": output["file_path"],
            "file_name": output["file_name"],
            "duration_secs": output["duration_secs"],
            "status": "completed",
            "project_title": request.project_title,
            "render_mode": request.render_mode,
            "creator_name": access_control.creator_name,
            "creator_role": access_control.creator_role,
            "access_level": access_control.access_level,
            "share_token": share_record["share_token"],
            "share_url": f"{get_api_origin()}/share/{share_record['share_token']}"
        }
        save_video_record(record)
        return VideoResponse(**record)
    except HTTPException:
        raise
    except Exception as error:
        print(f"Direct video generation failed: {str(error)}")
        raise HTTPException(status_code=500, detail=f"Video generation failed: {error}")


@router.get("/video/{video_id}")
async def get_video(video_id: str):
    library = load_json_list(LIBRARY_FILE)
    for item in library:
        if item.get("video_id") == video_id:
            return item

    video_path = os.path.join(config.OUTPUT_DIR, f"{video_id}.mp4")
    if not os.path.exists(video_path):
        raise HTTPException(status_code=404, detail=f"Video not found: {video_id}")

    return {
        "video_id": video_id,
        "file_path": f"outputs/{video_id}.mp4",
        "file_name": f"{video_id}.mp4",
        "status": "completed",
        "exists": True
    }


@router.get("/library")
async def get_video_library():
    videos = sorted(load_json_list(LIBRARY_FILE), key=lambda item: item.get("video_id", ""))
    return {"videos": videos}


@router.post("/share", response_model=ShareResponse)
async def create_share_link(request: ShareRequest):
    library = load_json_list(LIBRARY_FILE)
    video_record = next((item for item in library if item.get("video_id") == request.video_id), None)
    if not video_record:
        raise HTTPException(status_code=404, detail="Video not found")

    share_record = create_share_record(
        request.video_id,
        AccessControl(
            creator_name=request.creator_name,
            creator_role=request.creator_role,
            institution=request.institution,
            access_level=request.access_level,
            access_code=request.access_code,
            allow_download=request.allow_download
        )
    )

    video_record["share_token"] = share_record["share_token"]
    video_record["share_url"] = f"{get_api_origin()}/share/{share_record['share_token']}"
    video_record["access_level"] = share_record["access_level"]
    save_video_record(video_record)

    return ShareResponse(
        video_id=request.video_id,
        share_token=share_record["share_token"],
        share_url=video_record["share_url"],
        access_level=share_record["access_level"],
        access_code=share_record["access_code"] or None
    )


@router.get("/share/{share_token}")
async def get_shared_video(share_token: str, access_code: Optional[str] = None):
    shares = load_json_list(SHARES_FILE)
    share_record = next((item for item in shares if item.get("share_token") == share_token), None)
    if not share_record:
        raise HTTPException(status_code=404, detail="Share link not found")

    if share_record["access_level"] in ["private", "institution-only"]:
        expected_code = share_record.get("access_code") or ""
        if expected_code and access_code != expected_code:
            raise HTTPException(status_code=403, detail="Access code required or invalid")

    library = load_json_list(LIBRARY_FILE)
    video_record = next((item for item in library if item.get("video_id") == share_record["video_id"]), None)
    if not video_record:
        raise HTTPException(status_code=404, detail="Shared video record not found")

    return {
        "video_id": video_record["video_id"],
        "project_title": video_record.get("project_title"),
        "file_path": video_record["file_path"],
        "file_name": video_record["file_name"],
        "creator_name": share_record.get("creator_name"),
        "creator_role": share_record.get("creator_role"),
        "institution": share_record.get("institution"),
        "access_level": share_record.get("access_level"),
        "allow_download": share_record.get("allow_download", True)
    }


@router.get("/share/{share_token}/download")
async def download_shared_video(share_token: str, access_code: Optional[str] = None):
    shared_video = await get_shared_video(share_token, access_code)
    file_path = os.path.join(config.OUTPUT_DIR, shared_video["file_name"])
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="Video file not found")
    return FileResponse(path=file_path, filename=shared_video["file_name"], media_type="video/mp4")


@router.get("/video/{video_id}/status")
async def get_video_status(video_id: str):
    return {"video_id": video_id, "status": "completed", "progress": 100}


@router.delete("/video/{video_id}")
async def delete_video(video_id: str):
    library = load_json_list(LIBRARY_FILE)
    save_json_list(LIBRARY_FILE, [item for item in library if item.get("video_id") != video_id])

    video_path = os.path.join(config.OUTPUT_DIR, f"{video_id}.mp4")
    if os.path.exists(video_path):
        os.remove(video_path)

    return {"message": f"Video {video_id} deleted successfully"}


@router.get("/download/{filename}")
async def download_video(filename: str):
    file_path = os.path.join(config.OUTPUT_DIR, filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="Video file not found")
    return FileResponse(path=file_path, filename=filename, media_type="video/mp4")
